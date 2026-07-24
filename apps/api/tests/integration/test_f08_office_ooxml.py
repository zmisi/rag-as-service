"""F08 Office OOXML integration tests (upload + index + search)."""

from __future__ import annotations

import io
from uuid import uuid4

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from sqlalchemy import delete, select
from sqlalchemy.orm import Session

from rag_api.config import get_settings
from rag_api.db.models import Document as DocRow
from rag_api.db.models import DocumentChunk, DocumentFile, DocumentSection, IndexJob
from rag_api.indexing.embedding import HashingEmbedder
from rag_api.indexing.search import PgKnowledgeSearcher
from rag_api.indexing.worker import process_index_job
from tests.helpers import tenant_host_headers

HEADERS_A = tenant_host_headers("tenant-a")


def _docx_bytes(paragraph: str) -> bytes:
    doc = Document()
    if paragraph:
        doc.add_paragraph(paragraph)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes(slide_text: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 4000000, 1000000)
    box.text_frame.text = slide_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_bytes(cell: str, *, empty: bool = False) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    if not empty:
        ws["A1"] = cell
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture(autouse=True)
def wipe_docs(db: Session):
    db.execute(delete(DocumentChunk))
    db.execute(delete(DocumentSection))
    db.execute(delete(IndexJob))
    db.execute(delete(DocumentFile))
    db.execute(delete(DocRow))
    db.commit()


@pytest.fixture(autouse=True)
def sync_index(monkeypatch):
    get_settings.cache_clear()
    settings = get_settings()
    monkeypatch.setattr(settings, "index_sync_on_publish", True)
    yield
    get_settings.cache_clear()


def _create_doc(client) -> str:
    r = client.post("/v1/documents", headers=HEADERS_A)
    assert r.status_code == 201, r.text
    return r.json()["id"]


def _upload(client, doc_id: str, filename: str, data: bytes) -> None:
    r = client.post(
        f"/v1/documents/{doc_id}/files",
        headers=HEADERS_A,
        files={"file": (filename, io.BytesIO(data), "application/octet-stream")},
    )
    assert r.status_code == 201, r.text


def _prepare_meta(client, doc_id: str) -> None:
    r = client.patch(
        f"/v1/documents/{doc_id}",
        headers=HEADERS_A,
        json={"title": "office-doc", "tag": "faq"},
    )
    assert r.status_code == 200, r.text


def _submit_publish(client, doc_id: str) -> None:
    _prepare_meta(client, doc_id)
    r = client.post(f"/v1/documents/{doc_id}/submit-review", headers=HEADERS_A)
    assert r.status_code == 200, r.text
    r = client.post(f"/v1/documents/{doc_id}/publish", headers=HEADERS_A)
    assert r.status_code == 200, r.text


@pytest.mark.integration
def test_f08_t01_upload_ooxml_trio(client_a):
    doc_id = _create_doc(client_a)
    _upload(client_a, doc_id, "a.docx", _docx_bytes("p"))
    _upload(client_a, doc_id, "b.pptx", _pptx_bytes("s"))
    _upload(client_a, doc_id, "c.xlsx", _xlsx_bytes("c"))
    r = client_a.get(f"/v1/documents/{doc_id}", headers=HEADERS_A)
    assert r.status_code == 200
    names = {f["filename"] for f in r.json()["files"]}
    assert names == {"a.docx", "b.pptx", "c.xlsx"}


@pytest.mark.integration
def test_f08_t02_reject_legacy(client_a):
    doc_id = _create_doc(client_a)
    for name in ("x.doc", "x.ppt", "x.xls"):
        r = client_a.post(
            f"/v1/documents/{doc_id}/files",
            headers=HEADERS_A,
            files={"file": (name, io.BytesIO(b"OLE"), "application/octet-stream")},
        )
        assert r.status_code == 400, name
        detail = r.json()["detail"].lower()
        assert "docx" in detail and "pptx" in detail and "xlsx" in detail


@pytest.mark.integration
def test_f08_t09_t10_t11_magic_mismatch_upload(client_a):
    doc_id = _create_doc(client_a)
    cases = [
        ("bad.docx", b"not-zip"),
        ("bad.pdf", _docx_bytes("x")),  # zip pretending to be pdf
        ("bad.xlsx", _docx_bytes("x")),  # word/ zip as xlsx
    ]
    for name, data in cases:
        r = client_a.post(
            f"/v1/documents/{doc_id}/files",
            headers=HEADERS_A,
            files={"file": (name, io.BytesIO(data), "application/octet-stream")},
        )
        assert r.status_code == 400, name


@pytest.mark.integration
def test_f08_t03_docx_index_search(client_a, db: Session, tmp_path):
    phrase = f"DOCX_HIT_{uuid4().hex[:8]}"
    doc_id = _create_doc(client_a)
    _upload(client_a, doc_id, "m.docx", _docx_bytes(phrase))
    _submit_publish(client_a, doc_id)

    doc = db.get(DocRow, doc_id)
    assert doc is not None
    assert doc.index_status == "ready"

    # Confirm route via re-parse log path: sections non-empty
    n = db.scalar(
        select(DocumentSection).where(
            DocumentSection.doc_id == doc.doc_id,
            DocumentSection.is_latest.is_(True),
        )
    )
    assert n is not None
    assert phrase in (n.content or "")

    tenant_id = doc.tenant_id
    searcher = PgKnowledgeSearcher(lambda: db, embedder=HashingEmbedder())
    hits = searcher.search(tenant_id, phrase, top_k=5)
    assert any(phrase in (h.content or "") for h in hits)


@pytest.mark.integration
def test_f08_t04_pptx_index_search(client_a, db: Session):
    phrase = f"PPTX_HIT_{uuid4().hex[:8]}"
    doc_id = _create_doc(client_a)
    _upload(client_a, doc_id, "m.pptx", _pptx_bytes(phrase))
    _submit_publish(client_a, doc_id)
    doc = db.get(DocRow, doc_id)
    assert doc is not None and doc.index_status == "ready"
    searcher = PgKnowledgeSearcher(lambda: db, embedder=HashingEmbedder())
    hits = searcher.search(doc.tenant_id, phrase, top_k=5)
    assert any(phrase in (h.content or "") for h in hits)


@pytest.mark.integration
def test_f08_t05_xlsx_index_search(client_a, db: Session):
    phrase = f"XLSX_HIT_{uuid4().hex[:8]}"
    doc_id = _create_doc(client_a)
    _upload(client_a, doc_id, "m.xlsx", _xlsx_bytes(phrase))
    _submit_publish(client_a, doc_id)
    doc = db.get(DocRow, doc_id)
    assert doc is not None and doc.index_status == "ready"
    searcher = PgKnowledgeSearcher(lambda: db, embedder=HashingEmbedder())
    hits = searcher.search(doc.tenant_id, phrase, top_k=5)
    assert any(phrase in (h.content or "") for h in hits)


@pytest.mark.integration
def test_f08_t06_empty_ooxml_zero_chunks(client_a, db: Session):
    doc_id = _create_doc(client_a)
    _upload(client_a, doc_id, "empty.xlsx", _xlsx_bytes("", empty=True))
    _submit_publish(client_a, doc_id)
    doc = db.get(DocRow, doc_id)
    assert doc is not None
    assert doc.index_status == "ready"
    chunks = db.scalars(
        select(DocumentChunk).where(
            DocumentChunk.doc_id == doc.doc_id,
            DocumentChunk.is_latest.is_(True),
        )
    ).all()
    sections = db.scalars(
        select(DocumentSection).where(
            DocumentSection.doc_id == doc.doc_id,
            DocumentSection.is_latest.is_(True),
        )
    ).all()
    assert chunks == []
    assert sections == []


@pytest.mark.integration
def test_f08_t07_corrupt_ooxml_fails(client_a, db: Session, monkeypatch):
    """Extension+magic pass (word/ zip) but python-docx cannot open → job failed."""
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("word/document.xml", b"<not-a-real-docx/>")
    data = buf.getvalue()

    get_settings.cache_clear()
    settings = get_settings()
    monkeypatch.setattr(settings, "index_sync_on_publish", False)

    doc_id = _create_doc(client_a)
    _upload(client_a, doc_id, "bad.docx", data)
    _prepare_meta(client_a, doc_id)
    r = client_a.post(f"/v1/documents/{doc_id}/submit-review", headers=HEADERS_A)
    assert r.status_code == 200
    r = client_a.post(f"/v1/documents/{doc_id}/publish", headers=HEADERS_A)
    assert r.status_code == 200

    doc = db.get(DocRow, doc_id)
    assert doc is not None
    job = db.scalar(
        select(IndexJob)
        .where(IndexJob.doc_id == doc.doc_id)
        .order_by(IndexJob.create_at.desc())
    )
    assert job is not None
    from rag_api.services.storage_service import StorageService as SS

    with pytest.raises(Exception):
        process_index_job(db, job.id, embedder=HashingEmbedder(), storage=SS())
    db.refresh(job)
    db.refresh(doc)
    assert job.status == "failed"
    assert doc.index_status == "failed"
    latest = db.scalars(
        select(DocumentChunk).where(
            DocumentChunk.doc_id == doc.doc_id,
            DocumentChunk.is_latest.is_(True),
        )
    ).all()
    assert latest == []


# F08-T08: Phase 1 baseline rejected OOXML; with F08 enabled, T01 is authoritative.

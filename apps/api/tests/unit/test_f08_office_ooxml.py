"""F08 unit: magic-byte gates + lightweight Office → Markdown."""

from __future__ import annotations

import io
import zipfile

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from rag_api.domain.documents.constants import FILE_TYPE_MISMATCH_MESSAGE
from rag_api.domain.documents.file_type import FileTypeError, validate_file_type
from rag_api.indexing.office import (
    OfficeParseError,
    docx_to_markdown,
    pptx_to_markdown,
    xlsx_to_markdown,
)
from rag_api.indexing.parse import ParseError, RoutedDocumentParser


def _ooxml_zip(*names: str, extra: dict[str, bytes] | None = None) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        for name in names:
            zf.writestr(name, b"<xml/>")
        if extra:
            for k, v in extra.items():
                zf.writestr(k, v)
    return buf.getvalue()


def _real_docx(paragraph: str = "hello") -> bytes:
    doc = Document()
    doc.add_paragraph(paragraph)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _real_pptx(text: str = "slide text") -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = text if False else ""  # title may be None on blank
    # blank layout often has no title; add a text box
    box = slide.shapes.add_textbox(0, 0, 3000000, 1000000)
    box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _real_xlsx(cell: str = "cell") -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = cell
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_allows_xlsx_extension() -> None:
    from rag_api.domain.documents.constants import is_allowed_extension, is_legacy_extension

    assert is_allowed_extension("a.xlsx")
    assert not is_legacy_extension("a.xlsx")
    assert is_legacy_extension("a.xls")


def test_f08_t09_docx_name_non_zip() -> None:
    with pytest.raises(FileTypeError, match="match"):
        validate_file_type("x.docx", b"not a zip")


def test_f08_t10_pdf_name_zip_body() -> None:
    data = _ooxml_zip("word/document.xml")
    with pytest.raises(FileTypeError, match="match"):
        validate_file_type("x.pdf", data)


def test_f08_t11_xlsx_zip_without_xl() -> None:
    data = _ooxml_zip("word/document.xml")
    with pytest.raises(FileTypeError, match="match"):
        validate_file_type("book.xlsx", data)


def test_magic_ok_real_ooxml() -> None:
    assert validate_file_type("a.docx", _real_docx()) == ".docx"
    assert validate_file_type("a.pptx", _real_pptx()) == ".pptx"
    assert validate_file_type("a.xlsx", _real_xlsx()) == ".xlsx"
    assert validate_file_type("a.pdf", b"%PDF-1.4 hi") == ".pdf"
    assert validate_file_type("a.txt", b"hello") == ".txt"


def test_docx_to_markdown_phrase() -> None:
    md = docx_to_markdown(_real_docx("UNIQUE_DOCX_PHRASE_42"))
    assert "UNIQUE_DOCX_PHRASE_42" in md


def test_pptx_to_markdown_phrase() -> None:
    md = pptx_to_markdown(_real_pptx("UNIQUE_PPTX_PHRASE_42"))
    assert "UNIQUE_PPTX_PHRASE_42" in md


def test_xlsx_to_markdown_phrase() -> None:
    md = xlsx_to_markdown(_real_xlsx("UNIQUE_XLSX_PHRASE_42"))
    assert "UNIQUE_XLSX_PHRASE_42" in md
    assert "## Data" in md
    # Single cell is non-table prose, not a forced one-cell Markdown table.
    assert "| UNIQUE_XLSX_PHRASE_42 |" not in md


def test_f08_t13_xlsx_markdown_table_fill_down() -> None:
    """F08-T13: Markdown table + column fill-down for merged-cell empties."""
    from datetime import datetime

    from rag_api.indexing.sections import infer_chunk_type

    wb = Workbook()
    ws = wb.active
    ws.title = "产品信息"
    ws.append(["产品线", "产品型号", "价格(元)"])
    ws.append(["智能手机", "X100 Pro", 6999])
    ws.append([None, "X100", 4999])
    ws.append([None, "X90 Pro", 5999])
    ws.append(["平板电脑", "Pad Pro", 8999])
    ws2 = wb.create_sheet("混合数据类型")
    ws2.append(["文档名称", "发布日期"])
    ws2.append(["API开发指南", datetime(2026, 6, 15)])
    buf = io.BytesIO()
    wb.save(buf)
    md = xlsx_to_markdown(buf.getvalue())
    assert "## 产品信息" in md
    assert "| --- | --- | --- |" in md
    assert "| 智能手机 | X100 | 4999 |" in md
    assert "| 智能手机 | X90 Pro | 5999 |" in md
    assert "2026-06-15" in md
    assert "00:00:00" not in md
    body = md.split("## 产品信息", 1)[1].split("## ", 1)[0].strip()
    assert infer_chunk_type(body) == "table"


def test_f08_t14_xlsx_two_row_header_flatten() -> None:
    """F08-T14: sparse group row + dense leaf headers → leaf names as header."""
    wb = Workbook()
    ws = wb.active
    ws.title = "技术参数表"
    ws.append(["基础参数", None, "相机参数", None, None])
    ws.append(["产品型号", "处理器", "电池容量", "主摄", "超广角"])
    ws.append(["X100 Pro", "骁龙8 Gen 3", "5400mAh", "5000万", "5000万"])
    ws.append(["X100", "骁龙8 Gen 2", "5000mAh", "5000万", "1200万"])
    buf = io.BytesIO()
    wb.save(buf)
    md = xlsx_to_markdown(buf.getvalue())
    assert "| 产品型号 | 处理器 | 电池容量 | 主摄 | 超广角 |" in md
    assert "| --- | --- | --- | --- | --- |" in md
    assert "| X100 Pro |" in md
    data_lines = [
        ln
        for ln in md.splitlines()
        if ln.startswith("| ") and "---" not in ln and "产品型号 | 处理器" not in ln
    ]
    assert any("X100 Pro" in ln for ln in data_lines)
    assert not any("基础参数" in ln and "X100" not in ln for ln in data_lines)


def test_f08_t15_xlsx_multi_table_same_sheet() -> None:
    """F08-T15: blank row separates two tables on one sheet."""
    wb = Workbook()
    ws = wb.active
    ws.title = "混合"
    ws.append(["品类", "数量"])
    ws.append(["手机", 10])
    ws.append(["平板", 5])
    ws.append([None, None])
    ws.append(["城市", "仓位"])
    ws.append(["上海", "A1"])
    ws.append(["北京", "B2"])
    buf = io.BytesIO()
    wb.save(buf)
    md = xlsx_to_markdown(buf.getvalue())
    assert md.count("| --- | --- |") == 2
    assert "| 品类 | 数量 |" in md
    assert "| 城市 | 仓位 |" in md
    assert "| 手机 | 10 |" in md
    assert "| 上海 | A1 |" in md


def test_f08_t16_xlsx_cross_tab_to_long() -> None:
    """F08-T16: cross-tab matrix → long 行/列/值 table."""
    wb = Workbook()
    ws = wb.active
    ws.title = "交叉"
    ws.append([None, "Q1", "Q2", "Q3"])
    ws.append(["华东", 10, 20, 30])
    ws.append(["华北", 11, 21, 31])
    ws.append(["华南", 12, 22, 32])
    buf = io.BytesIO()
    wb.save(buf)
    md = xlsx_to_markdown(buf.getvalue())
    assert "| 行 | 列 | 值 |" in md
    assert "| 华东 | Q1 | 10 |" in md
    assert "| 华南 | Q3 | 32 |" in md
    assert "| 华东 | 10 | 20 | 30 |" not in md


def test_f08_t17_xlsx_prose_title_before_table() -> None:
    """F08-T17: single-cell title stays prose; table remains a table."""
    wb = Workbook()
    ws = wb.active
    ws.title = "说明表"
    ws.append(["2026年第二季度产品说明"])
    ws.append(["产品", "状态"])
    ws.append(["X100", "有货"])
    ws.append(["Pad", "缺货"])
    buf = io.BytesIO()
    wb.save(buf)
    md = xlsx_to_markdown(buf.getvalue())
    assert "2026年第二季度产品说明" in md
    assert "| 2026年第二季度产品说明 |" not in md
    assert "| 产品 | 状态 |" in md
    assert "| X100 | 有货 |" in md


def test_empty_docx_markdown() -> None:
    assert docx_to_markdown(_real_docx("")) == "" or docx_to_markdown(
        _real_docx("   ")
    ) == ""


def test_corrupt_docx_raises() -> None:
    # Valid ZIP+word/ magic but not a real docx package for python-docx
    data = _ooxml_zip("word/document.xml")
    with pytest.raises(OfficeParseError):
        docx_to_markdown(data)


def test_route_docx_not_docling() -> None:
    parser = RoutedDocumentParser()
    outcome = parser.parse_outcome("manual.docx", _real_docx("from-lite"))
    assert outcome.route == "docx"
    assert "from-lite" in outcome.text


def test_route_pptx_xlsx() -> None:
    parser = RoutedDocumentParser()
    assert parser.parse_outcome("a.pptx", _real_pptx("p")).route == "pptx"
    assert parser.parse_outcome("a.xlsx", _real_xlsx("x")).route == "xlsx"


def test_mismatch_raises_parse_error() -> None:
    parser = RoutedDocumentParser()
    with pytest.raises(ParseError):
        parser.parse_outcome("x.docx", b"plain")

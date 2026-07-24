"""Lightweight OOXML → Markdown (F08). No Docling."""

from __future__ import annotations

import datetime as dt
import io
import re
from pathlib import Path


class OfficeParseError(Exception):
    """Corrupt or unreadable OOXML package."""


def _heading_level_from_style(style_name: str | None) -> int | None:
    if not style_name:
        return None
    name = style_name.strip().lower()
    m = re.match(r"heading\s*(\d+)", name)
    if m:
        return max(1, min(int(m.group(1)), 6))
    if name in {"title", "subtitle"}:
        return 1 if name == "title" else 2
    return None


def _table_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [r + [""] * (width - len(r)) for r in rows]
    header = norm[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in norm[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _fill_down_empty_cells(rows: list[list[str]]) -> list[list[str]]:
    """Propagate previous non-empty cell down each column (merged-cell shape)."""
    if not rows:
        return rows
    width = max(len(r) for r in rows)
    last = [""] * width
    out: list[list[str]] = []
    for row in rows:
        padded = row + [""] * (width - len(row))
        filled: list[str] = []
        for i, cell in enumerate(padded):
            if cell:
                last[i] = cell
                filled.append(cell)
            else:
                filled.append(last[i])
        out.append(filled)
    return out


def _pad_row(row: list[str], width: int) -> list[str]:
    return row + [""] * (width - len(row))


def _trim_row_trailing(row: list[str]) -> list[str]:
    end = len(row)
    while end > 0 and not row[end - 1]:
        end -= 1
    return row[:end]


def _filled_ratio(row: list[str]) -> float:
    if not row:
        return 0.0
    return sum(1 for c in row if c) / len(row)


def _nonempty_count(row: list[str]) -> int:
    return sum(1 for c in row if c)


def _is_sparse_group_header(row: list[str]) -> bool:
    """First header row: some labels, many blanks (merged category cells)."""
    if not row or not any(row):
        return False
    ratio = _filled_ratio(row)
    return 0.0 < ratio < 0.6


def _is_dense_column_header(row: list[str]) -> bool:
    """Second header row: mostly short non-empty labels."""
    if not row:
        return False
    if _filled_ratio(row) < 0.8:
        return False
    filled = [c for c in row if c]
    return bool(filled) and all(len(c) <= 40 for c in filled)


def _looks_numeric(cell: str) -> bool:
    s = (cell or "").replace(",", "").replace("%", "").strip()
    if not s:
        return False
    try:
        float(s)
        return True
    except ValueError:
        return False


def _looks_like_table_block(block: list[list[str]]) -> bool:
    if len(block) < 2:
        return False
    width = max((len(r) for r in block), default=0)
    if width < 2:
        return False
    multi = sum(1 for r in block if _nonempty_count(r) >= 2)
    return multi >= max(2, (len(block) + 1) // 2)


def _split_sheet_blocks(rows: list[list[str] | None]) -> list[list[list[str]]]:
    """Split on fully empty rows into contiguous non-empty blocks."""
    blocks: list[list[list[str]]] = []
    current: list[list[str]] = []
    for row in rows:
        if row is None or not any(row):
            if current:
                blocks.append(current)
                current = []
            continue
        current.append(row)
    if current:
        blocks.append(current)
    return blocks


def _normalize_header_and_data(
    raw_rows: list[list[str]],
) -> tuple[list[str], list[list[str]]]:
    width = max(len(r) for r in raw_rows)
    rows = [_pad_row(r, width) for r in raw_rows]
    if (
        len(rows) >= 2
        and _is_sparse_group_header(rows[0])
        and _is_dense_column_header(rows[1])
    ):
        header = rows[1]
        data = _fill_down_empty_cells(rows[2:])
    else:
        header = rows[0]
        data = _fill_down_empty_cells(rows[1:])
    return header, data


def _is_cross_tab(header: list[str], data: list[list[str]]) -> bool:
    """Detect pivot/cross matrix: column dims on top, row dims on left, numeric body."""
    if len(data) < 2 or len(header) < 3:
        return False
    col_headers = header[1:]
    if sum(1 for c in col_headers if c) < 2:
        return False
    # Prefer classic empty top-left; also allow a short row-dimension label.
    if header[0] and len(header[0]) > 20:
        return False
    numeric = 0
    total = 0
    for row in data:
        if not row or not row[0] or _looks_numeric(row[0]):
            return False
        for cell in row[1 : len(header)]:
            if not cell:
                continue
            total += 1
            if _looks_numeric(cell):
                numeric += 1
    if total < 4:
        return False
    return (numeric / total) >= 0.7


def _cross_tab_to_long(
    header: list[str], data: list[list[str]]
) -> list[list[str]]:
    long_rows: list[list[str]] = [["行", "列", "值"]]
    col_names = header[1:]
    for row in data:
        row_label = row[0] if row else ""
        if not row_label:
            continue
        for idx, col_name in enumerate(col_names):
            if not col_name:
                continue
            val = row[idx + 1] if idx + 1 < len(row) else ""
            if not val:
                continue
            long_rows.append([row_label, col_name, val])
    return long_rows


def _matrix_to_markdown_table(raw_rows: list[list[str]]) -> str:
    """Build one Markdown table (with optional cross-tab → long)."""
    if not raw_rows:
        return ""
    header, data = _normalize_header_and_data(raw_rows)
    if data and _is_cross_tab(header, data):
        return _table_to_markdown(_cross_tab_to_long(header, data))
    return _table_to_markdown([header] + data)


def _block_to_markdown_parts(block: list[list[str]]) -> list[str]:
    """Emit prose and/or table markdown for one contiguous sheet block."""
    if not block:
        return []

    # Leading single-cell title/note before a table.
    if (
        len(block) >= 3
        and _nonempty_count(block[0]) == 1
        and _looks_like_table_block(block[1:])
    ):
        title = next(c for c in block[0] if c)
        rest = _block_to_markdown_parts(block[1:])
        return [title, *rest]

    if _looks_like_table_block(block):
        md = _matrix_to_markdown_table(block)
        return [md] if md else []

    # Non-table: join cells per row as prose lines.
    lines: list[str] = []
    for row in block:
        cells = [c for c in _trim_row_trailing(row) if c]
        if cells:
            lines.append(" ".join(cells))
    text = "\n".join(lines).strip()
    return [text] if text else []


def _sheet_rows_to_markdown_parts(rows: list[list[str] | None]) -> list[str]:
    parts: list[str] = []
    for block in _split_sheet_blocks(rows):
        parts.extend(_block_to_markdown_parts(block))
    return parts


def _cell_display(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, dt.datetime):
        if (
            value.hour == 0
            and value.minute == 0
            and value.second == 0
            and value.microsecond == 0
        ):
            return value.date().isoformat()
        return value.isoformat(sep=" ", timespec="seconds")
    if isinstance(value, dt.date):
        return value.isoformat()
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    text = str(value).strip()
    return text.replace("|", "\\|")


def docx_to_markdown(data: bytes) -> str:
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise OfficeParseError(
            "python-docx is not installed; pip install python-docx"
        ) from exc

    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise OfficeParseError(f"Failed to open .docx: {exc}") from exc

    parts: list[str] = []

    def emit_paragraph(p: Paragraph) -> None:
        text = (p.text or "").strip()
        if not text:
            return
        level = _heading_level_from_style(
            p.style.name if p.style is not None else None
        )
        if level:
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    def emit_table(table: Table) -> None:
        rows: list[list[str]] = []
        for row in table.rows:
            cells = [" ".join((c.text or "").split()) for c in row.cells]
            if any(cells):
                rows.append(cells)
        md = _table_to_markdown(rows)
        if md:
            parts.append(md)

    try:
        body = document.element.body
        for child in body.iterchildren():
            tag = child.tag
            if tag.endswith("}p"):
                emit_paragraph(Paragraph(child, document))
            elif tag.endswith("}tbl"):
                emit_table(Table(child, document))
    except Exception as exc:  # noqa: BLE001
        raise OfficeParseError(f"Failed to parse .docx content: {exc}") from exc

    return "\n\n".join(parts).strip()


def pptx_to_markdown(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise OfficeParseError(
            "python-pptx is not installed; pip install python-pptx"
        ) from exc

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise OfficeParseError(f"Failed to open .pptx: {exc}") from exc

    parts: list[str] = []
    try:
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            title = ""
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                chunk = (shape.text_frame.text or "").strip()
                if not chunk:
                    continue
                if not title:
                    title = chunk.splitlines()[0].strip()
                    rest = "\n".join(chunk.splitlines()[1:]).strip()
                    if rest:
                        texts.append(rest)
                else:
                    texts.append(chunk)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()

            heading = f"Slide {idx}"
            parts.append(f"## {heading}")
            body_bits: list[str] = []
            if title:
                body_bits.append(title)
            body_bits.extend(t for t in texts if t)
            if notes:
                body_bits.append(notes)
            if body_bits:
                parts.append("\n\n".join(body_bits))
    except Exception as exc:  # noqa: BLE001
        raise OfficeParseError(f"Failed to parse .pptx content: {exc}") from exc

    return "\n\n".join(parts).strip()


def xlsx_to_markdown(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise OfficeParseError(
            "openpyxl is not installed; pip install openpyxl"
        ) from exc

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        raise OfficeParseError(f"Failed to open .xlsx: {exc}") from exc

    parts: list[str] = []
    try:
        for sheet in wb.worksheets:
            name = sheet.title or "Sheet"
            parts.append(f"## {name}")
            # Keep empty rows as separators for multi-table sheets.
            raw_rows: list[list[str] | None] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [_cell_display(c) for c in row]
                if not any(cells):
                    raw_rows.append(None)
                else:
                    raw_rows.append(_trim_row_trailing(cells))
            for piece in _sheet_rows_to_markdown_parts(raw_rows):
                if piece:
                    parts.append(piece)
        wb.close()
    except Exception as exc:  # noqa: BLE001
        raise OfficeParseError(f"Failed to parse .xlsx content: {exc}") from exc

    return "\n\n".join(parts).strip()


def office_to_markdown(filename: str, data: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".docx":
        return docx_to_markdown(data)
    if suffix == ".pptx":
        return pptx_to_markdown(data)
    if suffix == ".xlsx":
        return xlsx_to_markdown(data)
    raise OfficeParseError(f"Not an OOXML office type: {suffix}")
